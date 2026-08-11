import re
from collections import defaultdict
from pptx import Presentation

from survey_parser.score_rules import METRIC_PATTERNS


def clean_text(text):
    if not text:
        return ""
    text = text.replace("\xa0", " ")
    text = text.replace("’", "'").replace("“", '"').replace("”", '"')
    text = re.sub(r"\s+", " ", text).strip()
    return text


def extract_text_blocks(ppt_path):
    prs = Presentation(ppt_path)
    blocks = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            # table cells
            if getattr(shape, "has_table", False):
                row_texts = []
                for row in shape.table.rows:
                    cell_row = []
                    for cell in row.cells:
                        t = clean_text(cell.text)
                        if t:
                            cell_row.append(t)
                            blocks.append({"slide": slide_idx, "text": t})
                    if cell_row:
                        row_joined = " | ".join(cell_row)
                        row_texts.append(row_joined)
                        blocks.append({"slide": slide_idx, "text": row_joined})

            # text frame
            if getattr(shape, "has_text_frame", False):
                paras = []
                for p in shape.text_frame.paragraphs:
                    t = clean_text(p.text)
                    if t:
                        paras.append(t)
                        blocks.append({"slide": slide_idx, "text": t})
                if paras:
                    blocks.append({"slide": slide_idx, "text": " | ".join(paras)})

            # fallback
            elif hasattr(shape, "text"):
                t = clean_text(shape.text)
                if t:
                    blocks.append({"slide": slide_idx, "text": t})

    return blocks


def _slide_text_map(blocks):
    slide_map = defaultdict(list)
    for b in blocks:
        slide_map[b["slide"]].append(b["text"])
    return {k: " | ".join(v) for k, v in slide_map.items()}


def _all_text(blocks):
    return "\n".join([b["text"] for b in blocks])


def _numbers_0_to_100(text):
    nums = [int(x) for x in re.findall(r"(?<!\d)(\d{1,3})(?!\d)", text)]
    return [n for n in nums if 0 <= n <= 100]


def _numbers_50_to_100(text):
    nums = _numbers_0_to_100(text)
    return [n for n in nums if 50 <= n <= 100]


def _extract_date(all_text):
    patterns = [
        r"((?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec|January|February|March|April|June|July|August|September|October|November|December)\s+\d{1,2},\s+\d{4})",
        r"((?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec|January|February|March|April|June|July|August|September|October|November|December)\s+'?\d{2})",
        r"((?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec|January|February|March|April|June|July|August|September|October|November|December)\s+\d{4})",
    ]
    for p in patterns:
        m = re.search(p, all_text, re.I)
        if m:
            return clean_text(m.group(1))
    return None


def _extract_manager(all_text):
    patterns = [
        r"manager[:\s]+([^\n|]+)",
        r"manager\s*-\s*([^\n|]+)",
    ]
    for p in patterns:
        m = re.search(p, all_text, re.I)
        if m:
            return clean_text(m.group(1))
    return None


def _extract_respondents(all_text):
    # e.g. 12 / 12 respondents
    m = re.search(r"(\d+)\s*/\s*(\d+)\s*respondents", all_text, re.I)
    if m:
        return int(m.group(1)), int(m.group(2))

    # backup: "12 respondents"
    m = re.search(r"(\d+)\s*respondents", all_text, re.I)
    if m:
        value = int(m.group(1))
        return value, value

    return None, None


def _extract_response_rate(all_text):
    # e.g. 12 (100%) Responded in Oct
    m = re.search(r"\(\s*(\d+)\s*%\s*\)\s*responded", all_text, re.I)
    if m:
        return int(m.group(1))

    # generic 100% response rate
    m = re.search(r"response\s*rate[^0-9]{0,12}(\d{1,3})\s*%", all_text, re.I)
    if m:
        return int(m.group(1))

    return None


def _extract_engagement_score(slide_map, all_text):
    # First try overview slide region
    for slide_no in sorted(slide_map.keys()):
        t = slide_map[slide_no]
        low = t.lower()
        if "engagement historical trend" in low or "survey overview" in low or "engagement1" in low:
            nums = _numbers_50_to_100(t)
            nums = [n for n in nums if n != 100]
            if nums:
                return nums[0]

    # Direct phrase matching
    patterns = [
        r"engagement[^0-9]{0,20}(\d{2,3})",
        r"overall\s+engagement[^0-9]{0,20}(\d{2,3})",
    ]
    for p in patterns:
        m = re.search(p, all_text, re.I)
        if m:
            value = int(m.group(1))
            if 50 <= value <= 100:
                return value

    return None


def _extract_company_score(slide_map, all_text):
    for slide_no in sorted(slide_map.keys()):
        t = slide_map[slide_no]
        m = re.search(r"company[^0-9]{0,12}(\d{2,3})", t, re.I)
        if m:
            value = int(m.group(1))
            if 50 <= value <= 100:
                return value

    m = re.search(r"company[^0-9]{0,12}(\d{2,3})", all_text, re.I)
    if m:
        value = int(m.group(1))
        if 50 <= value <= 100:
            return value

    return None


def _extract_score_near_alias(text, alias):
    """
    Find metric score around alias.
    We prefer 50-100 values and ignore very small change values like +2, +3.
    """
    alias = re.escape(alias)

    patterns = [
        rf"{alias}[^0-9]{{0,35}}(\d{{2,3}})",
        rf"(\d{{2,3}})[^0-9]{{0,35}}{alias}",
    ]

    for p in patterns:
        m = re.search(p, text, re.I)
        if m:
            value = int(m.group(1))
            if 50 <= value <= 100:
                return value

    nums = _numbers_50_to_100(text)
    if nums:
        # Prefer highest valid favorability score in the local snippet
        return max(nums)

    return None


def _match_metrics_in_text(text):
    found = []
    text_low = text.lower()

    for canonical, aliases in METRIC_PATTERNS.items():
        for alias in aliases:
            if alias in text_low:
                score = _extract_score_near_alias(text_low, alias.lower())
                if score is not None:
                    found.append({"name": canonical, "score": score})
                    break
    return found


def _extract_metrics_from_blocks(blocks):
    metrics = []

    # 1) block-level extraction
    for b in blocks:
        text = b["text"]
        metrics.extend(_match_metrics_in_text(text))

    # 2) slide-level extraction for fragmented PPT layouts
    slide_map = _slide_text_map(blocks)
    for slide_no, slide_text in slide_map.items():
        metrics.extend(_match_metrics_in_text(slide_text))

    # 3) score ranking / deduplicate
    grouped = defaultdict(list)
    for m in metrics:
        grouped[m["name"]].append(m["score"])

    final_metrics = []
    for name, scores in grouped.items():
        # choose the most plausible score:
        # - use most frequent if repeated
        # - otherwise highest score
        freq = defaultdict(int)
        for s in scores:
            freq[s] += 1
        best_score = sorted(freq.items(), key=lambda x: (x[1], x[0]), reverse=True)[0][0]
        final_metrics.append({"name": name, "score": best_score})

    final_metrics = sorted(final_metrics, key=lambda x: x["score"], reverse=True)
    return final_metrics


def _derive_priority_themes(metrics):
    """
    Used later for page 2 summary logic.
    """
    if not metrics:
        return []

    low_metrics = sorted(metrics, key=lambda x: x["score"])[:5]
    themes = []
    for m in low_metrics:
        themes.append(m["name"])
    return themes[:3]


def extract_overview(blocks):
    slide_map = _slide_text_map(blocks)
    all_text = _all_text(blocks)

    respondents, total_respondents = _extract_respondents(all_text)
    response_rate = _extract_response_rate(all_text)

    if response_rate is None and respondents and total_respondents:
        response_rate = round((respondents / total_respondents) * 100)

    overview = {
        "survey_date": _extract_date(all_text),
        "manager_or_team": _extract_manager(all_text),
        "engagement_score": _extract_engagement_score(slide_map, all_text),
        "company_score": _extract_company_score(slide_map, all_text),
        "respondents": respondents,
        "total_respondents": total_respondents,
        "response_rate": response_rate,
    }

    return overview


def parse_survey_ppt(ppt_path):
    blocks = extract_text_blocks(ppt_path)
    overview = extract_overview(blocks)
    metrics = _extract_metrics_from_blocks(blocks)
    priority_themes = _derive_priority_themes(metrics)

    return {
        "overview": overview,
        "metrics": metrics,
        "priority_themes": priority_themes,
        "text_blocks": blocks,
    }
