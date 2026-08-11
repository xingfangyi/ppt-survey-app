from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from ppt_generator.template_config import ABB_RED, ABB_LILAC, BLACK


def _rgb(hex_color):
    hex_color = hex_color.replace("#", "")
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16),
    )


def _set_shape_text(shape, lines):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True

    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = line.get("align", PP_ALIGN.LEFT)
        p.level = line.get("level", 0)

        run = p.add_run()
        run.text = line["text"]
        run.font.size = Pt(line.get("size", 14))
        run.font.bold = line.get("bold", False)
        run.font.color.rgb = _rgb(line.get("color", BLACK))


def _safe_title(slide, text):
    if slide.shapes.title is not None:
        slide.shapes.title.text = text


def _safe_subtitle(slide, text):
    for shape in slide.shapes:
        if getattr(shape, "is_placeholder", False) and hasattr(shape, "text_frame"):
            if shape != slide.shapes.title:
                try:
                    shape.text = text
                    return
                except Exception:
                    pass


def _metric_lines(metrics, limit=None):
    if not metrics:
        return ["No metric detected."]
    rows = metrics[:limit] if limit else metrics
    return [f"{m['name']} — {m['score']}" for m in rows]


def _overview_bullets(overview):
    bullets = []

    if overview.get("manager_or_team"):
        bullets.append(f"Team / Manager: {overview['manager_or_team']}")
    if overview.get("survey_date"):
        bullets.append(f"Survey date: {overview['survey_date']}")
    if overview.get("engagement_score") is not None:
        bullets.append(f"Engagement score: {overview['engagement_score']}")
    if overview.get("company_score") is not None:
        bullets.append(f"Company benchmark: {overview['company_score']}")
    if overview.get("respondents") is not None and overview.get("total_respondents") is not None:
        bullets.append(
            f"Respondents: {overview['respondents']} / {overview['total_respondents']}"
        )
    if overview.get("response_rate") is not None:
        bullets.append(f"Response rate: {overview['response_rate']}%")

    if not bullets:
        bullets.append("No overview field detected from source PPT.")

    return bullets


def _replace_if_possible(slide, shape_idx, heading, bullets):
    if shape_idx < len(slide.shapes):
        shape = slide.shapes[shape_idx]
        if hasattr(shape, "text_frame"):
            lines = [{"text": heading, "size": 18, "bold": True, "color": ABB_RED}]
            for b in bullets:
                lines.append({"text": b, "size": 14, "color": BLACK})
            _set_shape_text(shape, lines)
            return True
    return False


def _remove_shape(shape):
    sp = shape._element
    sp.getparent().remove(sp)


def _remove_pictures_on_slide(slide):
    """
    Remove picture objects on a slide.
    This is used on template content slides where the original
    content is stored as screenshots/images.
    """
    pictures = [shape for shape in slide.shapes if shape.shape_type == 13]
    for shape in pictures:
        _remove_shape(shape)


def _add_content_panel(slide, left, top, width, height, heading, bullets, heading_color=ABB_RED):
    """
    Create a white content panel and place text inside it.
    """
    rect = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = _rgb("FFFFFF")
    rect.line.color.rgb = _rgb("E6E6E6")
    rect.line.width = Pt(1.0)

    tf = rect.text_frame
    tf.clear()
    tf.word_wrap = True

    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = heading
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = _rgb(heading_color)

    for b in bullets:
        p = tf.add_paragraph()
        p.level = 0
        r = p.add_run()
        r.text = f"• {b}"
        r.font.size = Pt(13)
        r.font.color.rgb = _rgb(BLACK)


def _add_two_column_panels(slide, left_heading, left_bullets, right_heading, right_bullets):
    # Reuse the visual area previously occupied by template pictures
    top = 1.75
    height = 3.3
    _add_content_panel(
        slide,
        Inches(0.35), Inches(top), Inches(4.75), Inches(height),
        left_heading, left_bullets, ABB_RED
    )
    _add_content_panel(
        slide,
        Inches(5.15), Inches(top), Inches(4.75), Inches(height),
        right_heading, right_bullets, ABB_LILAC
    )


def _add_single_wide_panel(slide, heading, bullets):
    _add_content_panel(
        slide,
        Inches(0.35), Inches(1.55), Inches(9.55), Inches(4.55),
        heading, bullets, ABB_RED
    )


def generate_ppt(template_path, output_path, data):
    template_path = Path(template_path)
    prs = Presentation(str(template_path))

    overview = data["overview"]
    strengths = data["strengths"]
    opportunities = data["opportunities"]
    bottom10 = data["bottom10"]
    top3 = data["top3"]
    bottom3 = data["bottom3"]
    actions = data["actions"]
    key_messages = data["key_messages"]

    # Slide 1 - cover
    if len(prs.slides) >= 1:
        slide = prs.slides[0]
        team_name = overview.get("manager_or_team") or "Team Survey"
        _safe_title(slide, "Development Actions – Follow up in 2026")
        _safe_subtitle(slide, f"Based on Engagement Survey results, {team_name}")

    # Slide 2 - overview
    if len(prs.slides) >= 2:
        slide = prs.slides[1]
        _safe_title(slide, "Team Overview")
        _remove_pictures_on_slide(slide)
        _add_two_column_panels(
            slide,
            "Survey Snapshot",
            _overview_bullets(overview),
            "Key Messages",
            key_messages,
        )

    # Slide 3 - strengths and opportunities
    if len(prs.slides) >= 3:
        slide = prs.slides[2]
        _safe_title(slide, "Strengths and Opportunities")
        _remove_pictures_on_slide(slide)
        _add_two_column_panels(
            slide,
            "Top Strengths",
            _metric_lines(strengths, 5),
            "Top Opportunities",
            _metric_lines(opportunities, 5),
        )

    # Slide 4 - bottom 10
    if len(prs.slides) >= 4:
        slide = prs.slides[3]
        _safe_title(slide, "Bottom 10 Scores")
        _remove_pictures_on_slide(slide)
        _add_two_column_panels(
            slide,
            "Bottom 10 (1–5)",
            _metric_lines(bottom10[:5], 5),
            "Bottom 10 (6–10)",
            _metric_lines(bottom10[5:10], 5),
        )

    # Slide 5 - summary detail
    if len(prs.slides) >= 5:
        slide = prs.slides[4]
        _safe_title(slide, "Survey Summary")
        _remove_pictures_on_slide(slide)
        _add_two_column_panels(
            slide,
            "Top 3 Scores",
            _metric_lines(top3, 3),
            "Management Notes",
            key_messages,
        )

    # Slide 6 - focus areas
    if len(prs.slides) >= 6:
        slide = prs.slides[5]
        _safe_title(slide, "Focus Areas")
        _remove_pictures_on_slide(slide)
        left_bullets = [f"{x['name']} — {x['score']}" for x in opportunities[:5]] or ["No focus area detected."]
        right_bullets = [f"{a['theme']}: {a['quote']}" for a in actions]
        _add_two_column_panels(
            slide,
            "Improvement Topics",
            left_bullets,
            "Action Themes",
            right_bullets,
        )

    # Slide 7 - top / bottom 3
    if len(prs.slides) >= 7:
        slide = prs.slides[6]
        _safe_title(slide, "Top / Bottom 3 Scores")
        _remove_pictures_on_slide(slide)
        _add_two_column_panels(
            slide,
            "Top 3",
            _metric_lines(top3, 3),
            "Bottom 3",
            _metric_lines(bottom3, 3),
        )

    # Slide 8 - development actions
    if len(prs.slides) >= 8:
        slide = prs.slides[7]
        _safe_title(slide, "Development Actions – Follow Up in 2026")

        action1 = actions[0] if len(actions) >= 1 else {
            "theme": "Belonging",
            "quote": "I feel a sense of belonging at ABB.",
            "bullets": ["Create regular team connection activities."]
        }
        action2 = actions[1] if len(actions) >= 2 else {
            "theme": "Work-Life Balance",
            "quote": "I am able to successfully balance my work and personal life.",
            "bullets": ["Review workload and reduce overtime peaks."]
        }

        text1 = [f'{action1["theme"]}: "{action1["quote"]}"', ""] + action1["bullets"]
        text2 = [f'{action2["theme"]}: "{action2["quote"]}"', ""] + action2["bullets"]

        ok1 = _replace_if_possible(slide, 3, text1[0], text1[2:])
        ok2 = _replace_if_possible(slide, 5, text2[0], text2[2:])

        if not ok1:
            _add_content_panel(
                slide, Inches(0.35), Inches(2.3), Inches(4.7), Inches(2.9),
                text1[0], text1[2:], ABB_RED
            )
        if not ok2:
            _add_content_panel(
                slide, Inches(5.1), Inches(2.3), Inches(4.7), Inches(2.9),
                text2[0], text2[2:], ABB_LILAC
            )

        _replace_if_possible(slide, 8, "Professional Leadership Development", [
            "Managers review actions monthly.",
            "Track progress and close the loop with the team.",
            "Use survey insights in team discussions and planning.",
        ])
        _replace_if_possible(slide, 9, f"Improvement topic: {action1['theme']}", [])
        _replace_if_possible(slide, 10, f"Improvement topic: {action2['theme']}", [])

    # Slide 9 - closing date
    if len(prs.slides) >= 9:
        slide = prs.slides[8]
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                txt = getattr(shape, "text", "")
                if "2026" in txt or "Slide 9" in txt or txt.strip():
                    try:
                        shape.text = datetime.now().strftime("%B %d, %Y")
                        break
                    except Exception:
                        pass

    prs.save(str(output_path))
